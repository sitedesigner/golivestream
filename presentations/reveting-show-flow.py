#!/usr/bin/env python3
"""
Reveting Show Operations — End-to-End Presentation Generator
Generates a 13-slide PowerPoint deck for show production.

Usage:
  python3 presentations/reveting-show-flow.py
  python3 presentations/reveting-show-flow.py --output custom-name.pptx
  python3 presentations/reveting-show-flow.py --show "The David Daily Show" --guest "Daniel Burrus"
"""

import argparse
import os
import sys
from datetime import datetime

try:
    import ctypes
    ctypes.CDLL("/usr/lib/libexpat.1.dylib")
except (OSError, AttributeError):
    pass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("Installing python-pptx...")
    os.system(f"{sys.executable} -m pip install --break-system-packages python-pptx 2>/dev/null")
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Brand Colors ──────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x3E)
GOLD    = RGBColor(0xF5, 0xA6, 0x23)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF0, 0xF4, 0xFF)
GREEN   = RGBColor(0x1E, 0xA5, 0x6B)
RED     = RGBColor(0xD9, 0x3B, 0x3B)
MIDBLUE = RGBColor(0x1A, 0x4B, 0x9F)
GREY    = RGBColor(0x8A, 0x9B, 0xB5)
DARK    = RGBColor(0x0A, 0x12, 0x28)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
BLANK_LAYOUT = 6  # blank layout index

prs = None


def add_slide(title_text="", subtitle_text=""):
    """Add a new slide with optional title."""
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    # Header bar
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_W, Inches(1.0)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()

    if title_text:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.color.rgb = WHITE
        p.font.bold = True

    if subtitle_text:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(14)
        p.font.color.rgb = GOLD

    return slide


def add_text(slide, left, top, width, height, text, font_size=14,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=12, color=WHITE):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return txBox


def add_box(slide, left, top, width, height, fill_color, text="", text_color=WHITE):
    """Add a colored rectangle with optional text."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = text_color
        p.font.bold = True
    return shape


def add_step_boxes(slide, steps, y_start=Inches(1.5)):
    """Add numbered step boxes in a row."""
    n = len(steps)
    box_w = Inches(11.0 / n) - Inches(0.2)
    start_x = Inches(0.6)

    for i, (title, desc) in enumerate(steps):
        x = start_x + i * (box_w + Inches(0.2))

        # Number circle
        circle = slide.shapes.add_shape(9, x + Inches(0.1), y_start, Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = GOLD
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = NAVY
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Box
        box = add_box(slide, x, y_start + Inches(0.5), box_w, Inches(1.5),
                      MIDBLUE, title + "\n" + desc)
        return


def make_title_slide():
    """Slide 1: Title"""
    slide = add_slide("REVETING SHOW OPERATIONS", "End-to-End Production Flow")

    # 5 pillars
    pillars = ["SHOW SETUP", "GUEST BOOKING", "EMAIL SEQUENCE", "LIVE SHOW", "POST-PROD"]
    colors = [GOLD, GREEN, MIDBLUE, RED, GOLD]
    box_w = Inches(2.2)
    start_x = Inches(0.6)

    for i, (pillar, color) in enumerate(zip(pillars, colors)):
        x = start_x + i * (box_w + Inches(0.15))
        add_box(slide, x, Inches(2.5), box_w, Inches(1.2), color, pillar)

    # Source credit
    add_text(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
             "Based on Jessie Lizak (Reveting) LinkedIn Live & Livestream Content Engine SOP",
             font_size=10, color=GREY)


def make_overview_slide():
    """Slide 2: Complete Show Flow Overview"""
    slide = add_slide("COMPLETE SHOW FLOW", "4-Column Production Pipeline")

    lanes = [
        ("TIER 1: SHOW SETUP", ["Show Identity", "Team Roster", "Channel Stack",
                                  "Podcast Links", "Booking Docs", "Guest Rules"], GOLD),
        ("EPISODE PREP", ["Guest Booking", "Fit Gate", "Tier 2 Variables",
                           "Topic Editing", "Calendar Stage 1", "T-14 Activation"], GREEN),
        ("PRE-SHOW EMAILS", ["Email 1 (Event Published)", "Email 2 (T-7)",
                              "Email 3 (Day Before)", "Email 4 (Day Of)",
                              "Calendar Stage 2", "StreamYard Guest Link"], MIDBLUE),
        ("LIVE & POST-SHOW", ["Pre-show Tech Check", "Go Live (60 min)",
                               "Email 5 (Within 1h)", "Clips Delivered",
                               "Email 6 (24-36h)", "Podcast Upload"], RED),
    ]

    for i, (title, items, color) in enumerate(lanes):
        x = Inches(0.3) + i * Inches(3.2)
        add_box(slide, x, Inches(1.3), Inches(3.0), Inches(5.5), DARK)
        add_text(slide, x + Inches(0.1), Inches(1.4), Inches(2.8), Inches(0.4),
                 title, font_size=11, color=color, bold=True)
        for j, item in enumerate(items):
            add_text(slide, x + Inches(0.2), Inches(1.9) + j * Inches(0.65),
                     Inches(2.6), Inches(0.6), "• " + item, font_size=9, color=WHITE)


def make_tier1_slide():
    """Slide 3: Tier 1 Show Setup"""
    slide = add_slide("PHASE 1: NEW SHOW SETUP", "Tier 1 — Complete once at show launch")

    groups = [
        ("Show Identity", "[SHOW_NAME]\n[SHOW_EMAIL]\n[SHOW_TIMEZONE] → Eastern\n[SHOW_DAY_OF_WEEK]\n[SHOW_RECURRING_TIME_ET]"),
        ("Team", "[HOST_NAME]\n[HOST_LINKEDIN_URL]\n[PRODUCER_NAME]\n[PA_NAME]\n[SIGNATURE_BLOCK]"),
        ("Channel Stack", "⚠ Confirm against Total\nDeliverables Doc\nNEVER guess\nLinkedIn / YouTube / FB\nTwitch / Instagram"),
        ("Podcast Links", "Spotify\nApple Podcasts\nYouTube Channel\nShow Website\nPast Episodes"),
        ("Booking & Docs", "Booking Form\nTotal Deliverables Doc\nStreamYard Account\nClient Outline\nGoogle Drive Root\nLinkedIn Walkthrough"),
        ("Guest Rules", "5,000+ followers preferred\n1,000-4,999: confirm\nBelow 1,000: NO-GO\nEscalation → WhatsApp Jessie"),
    ]

    for i, (title, content) in enumerate(groups):
        col = i % 3
        row = i // 3
        x = Inches(0.4) + col * Inches(4.2)
        y = Inches(1.5) + row * Inches(2.8)
        add_box(slide, x, y, Inches(4.0), Inches(2.5), MIDBLUE)
        add_text(slide, x + Inches(0.1), y + Inches(0.1), Inches(3.8), Inches(0.3),
                 title, font_size=12, color=GOLD, bold=True)
        add_text(slide, x + Inches(0.1), y + Inches(0.4), Inches(3.8), Inches(1.8),
                 content, font_size=9, color=WHITE)


def make_guest_booking_slide():
    """Slide 4: Guest Booking & Fit Gate"""
    slide = add_slide("PHASE 2: GUEST BOOKING & FIT GATE", "Episode intake from app.reveting.com")

    # Left: Booking steps
    add_box(slide, Inches(0.4), Inches(1.3), Inches(5.5), Inches(5.5), DARK)
    add_text(slide, Inches(0.6), Inches(1.4), Inches(5), Inches(0.4),
             "BOOKING INTAKE STEPS", font_size=14, color=GOLD, bold=True)

    steps = [
        "1. Find Booking → Calendars → Appointment List View",
        "2. Pull Form Submission → 3 dots → View Details",
        "3. Complete Tier 2 Guest Info",
        "4. Check LinkedIn Followers (manual lookup)",
        "5. Edit Topics → 3-7 words each, aligned to strategy",
        "6. Set Episode Details (number, date, title)",
        "7. Calendar Description Stage 1 (immediate)",
    ]
    for j, step in enumerate(steps):
        add_text(slide, Inches(0.7), Inches(1.9) + j * Inches(0.55),
                 Inches(5), Inches(0.5), step, font_size=10, color=WHITE)

    # Right: Fit Gate
    add_box(slide, Inches(6.5), Inches(1.3), Inches(6.3), Inches(5.5), DARK)
    add_text(slide, Inches(6.7), Inches(1.4), Inches(6), Inches(0.4),
             "FIT GATE DECISION TREE", font_size=14, color=GOLD, bold=True)

    add_box(slide, Inches(6.7), Inches(2.0), Inches(5.8), Inches(0.8), GREEN)
    add_text(slide, Inches(6.9), Inches(2.1), Inches(5.4), Inches(0.6),
             "≥ 5,000 followers → ✓ FIT → Proceed", font_size=12, color=WHITE, bold=True)

    add_box(slide, Inches(6.7), Inches(3.0), Inches(5.8), Inches(0.8), GOLD)
    add_text(slide, Inches(6.9), Inches(3.1), Inches(5.4), Inches(0.6),
             "1,000-4,999 → ⚠ Confirm with Prod + Marketing Lead", font_size=11, color=NAVY, bold=True)

    add_box(slide, Inches(6.7), Inches(4.0), Inches(5.8), Inches(0.8), RED)
    add_text(slide, Inches(6.9), Inches(4.1), Inches(5.4), Inches(0.6),
             "< 1,000 → ✗ NOT A FIT → Send rejection email → STOP", font_size=11, color=WHITE, bold=True)

    add_text(slide, Inches(6.7), Inches(5.2), Inches(5.8), Inches(0.5),
             "⚠ Never email raw submitted topics. Edit to fit show strategy first.",
             font_size=10, color=GOLD, bold=True)


def make_t14_slide():
    """Slide 5: T-14 Activation & Calendar Stage 1"""
    slide = add_slide("PHASE 3: T-14 ACTIVATION", "14 days before show (or immediate if < 14 days)")

    # T-14 Email
    add_box(slide, Inches(0.4), Inches(1.3), Inches(6.0), Inches(3.5), DARK)
    add_text(slide, Inches(0.6), Inches(1.4), Inches(5.6), Inches(0.4),
             "T-14 ACTIVATION EMAIL", font_size=14, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(0.7), Inches(1.9), Inches(5.4), Inches(2.5),
                    ["Subject: Confirming your appearance on [SHOW_NAME]",
                     "Confirm date + time (all 3 time zones)",
                     "Request LinkedIn connection with [LINKEDIN_CHANNEL_OWNER]",
                     "Link to past episodes",
                     "Explain topic editing process",
                     "[SIGNATURE_BLOCK]"], font_size=10)

    # Calendar Stage 1
    add_box(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(3.5), DARK)
    add_text(slide, Inches(7.0), Inches(1.4), Inches(5.6), Inches(0.4),
             "CALENDAR STAGE 1 (Set Immediately)", font_size=14, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(7.1), Inches(1.9), Inches(5.4), Inches(2.5),
                    ["Event: [SHOW_NAME] with [GUEST_FIRST_NAME] [GUEST_LAST_NAME]",
                     "Location: StreamYard URL TBD",
                     "\"Thank you for submitting your topics. We will format them...\"",
                     "Login 15 min early at [PRESHOW_TIME_ET/CT/PT]",
                     "Go live: [GOLIVE_TIME_ET/CT/PT]",
                     "⚠ Preserve original booking form submission at bottom"],
                    font_size=10)


def make_email1_slide():
    """Slide 6: Email 1 — Event Published"""
    slide = add_slide("EMAIL 1: EVENT PUBLISHED", "Most variable-heavy email — all Tier 2 details required")

    add_box(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.5), DARK)
    add_text(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
             "TRIGGER: Production + Marketing Lead has published event on all platforms",
             font_size=12, color=GOLD, bold=True)

    add_text(slide, Inches(0.6), Inches(1.9), Inches(12), Inches(0.3),
             "REQUIRED VARIABLES:", font_size=11, color=RED, bold=True)

    vars_text = ("[EPISODE_DATE] [EPISODE_TITLE] [TOPIC_1-4] "
                 "[GOLIVE_TIME_ET/CT/PT] [STREAMYARD_GUEST_LINK] "
                 "[LINKEDIN_EVENT_URL] [SHOW_WEBSITE_LINK] "
                 "[HOST_NAME] [HOST_LINKEDIN_URL] [SIGNATURE_BLOCK]")
    add_text(slide, Inches(0.6), Inches(2.2), Inches(12), Inches(0.4),
             vars_text, font_size=10, color=LIGHT)

    add_text(slide, Inches(0.6), Inches(2.8), Inches(12), Inches(0.3),
             "CONTENT:", font_size=11, color=GOLD, bold=True)

    add_bullet_list(slide, Inches(0.7), Inches(3.2), Inches(11.5), Inches(3.0),
                    ["All episode details (date, title, topics, timing in 3 time zones)",
                     "StreamYard guest link (fresh!)",
                     "Promo links (LinkedIn event, show website, podcast links)",
                     "LinkedIn invite walkthrough instructions",
                     "Signature block"], font_size=10)


def make_preshow_slide():
    """Slide 7: Pre-Show Email Cadence"""
    slide = add_slide("PRE-SHOW EMAIL CADENCE", "Emails 2, 3, 4 — Building momentum")

    emails = [
        ("EMAIL 2 — T-7 DAYS", "Remind guest to invite LinkedIn connections",
         ["Briefing call booking link", "Show format overview", "Tech check preview"],
         GREEN),
        ("EMAIL 3 — DAY BEFORE", "Send StreamYard link again + hype",
         ["StreamYard link (all 3 time zones)", "Episode title + topics",
          "Equipment checklist reminder"], MIDBLUE),
        ("EMAIL 4 — MORNING OF", "Final reminder 2-4 hours before",
         ["StreamYard link (all 3 time zones)", "Go live time",
          "Reply READY or RESCHEDULE"], GOLD),
    ]

    for i, (title, purpose, bullets, color) in enumerate(emails):
        x = Inches(0.4) + i * Inches(4.2)
        add_box(slide, x, Inches(1.3), Inches(4.0), Inches(5.5), DARK)
        add_text(slide, x + Inches(0.1), Inches(1.4), Inches(3.8), Inches(0.3),
                 title, font_size=12, color=color, bold=True)
        add_text(slide, x + Inches(0.1), Inches(1.7), Inches(3.8), Inches(0.3),
                 purpose, font_size=9, color=GREY)
        for j, b in enumerate(bullets):
            add_text(slide, x + Inches(0.2), Inches(2.1) + j * Inches(0.5),
                     Inches(3.6), Inches(0.4), "• " + b, font_size=9, color=WHITE)


def make_streamyard_slide():
    """Slide 8: Calendar Stage 2 & StreamYard Prep"""
    slide = add_slide("CALENDAR STAGE 2 & STREAMYARD PREP", "When event is published across platforms")

    # Calendar Stage 2
    add_box(slide, Inches(0.4), Inches(1.3), Inches(6.0), Inches(3.5), DARK)
    add_text(slide, Inches(0.6), Inches(1.4), Inches(5.6), Inches(0.4),
             "CALENDAR STAGE 2 (Event Published)", font_size=14, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(0.7), Inches(1.9), Inches(5.4), Inches(2.5),
                    ["Finalized title with topics",
                     "All platform links (LinkedIn, YouTube, Facebook, etc.)",
                     "StreamYard guest link",
                     "Podcast links",
                     "[SIGNATURE_BLOCK]"], font_size=10)

    # StreamYard Prep
    add_box(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(3.5), DARK)
    add_text(slide, Inches(7.0), Inches(1.4), Inches(5.6), Inches(0.4),
             "STREAMYARD PREP CHECKLIST", font_size=14, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(7.1), Inches(1.9), Inches(5.4), Inches(2.5),
                    ["6 scenes built (Holding, HostSolo, HostGuest, ScreenShare, Quote, EndCard)",
                     "Lower thirds with guest name/title/company",
                     "Logo bug at 30-40% opacity",
                     "Guest invite link created",
                     "Multi-destination streaming tested"], font_size=10)


def make_live_slide():
    """Slide 9: Live Show Run-of-Show"""
    slide = add_slide("LIVE SHOW RUN-OF-SHOW", "45-Minute Interview Template")

    times = [
        ("T-5", "Holding screen live", "Holding"),
        ("T+0", "Host intro, guest intro", "Host Solo → Host+Guest"),
        ("T+3", "Guest bio, origin story", "Host+Guest"),
        ("T+12", "Chapter 1: Topic deep-dive", "Host+Guest"),
        ("T+24", "Chapter 2: Second topic", "Host+Guest"),
        ("T+36", "Chapter 3: Final topic + Q&A", "Host+Guest"),
        ("T+44", "CTA, next episode", "Host Solo"),
        ("T+45", "End card", "End Card"),
    ]

    for i, (time, action, scene) in enumerate(times):
        x = Inches(0.4) + i * Inches(1.55)
        add_box(slide, x, Inches(1.5), Inches(1.45), Inches(0.6), MIDBLUE,
                 time + "\n" + action, NAVY)
        add_text(slide, x, Inches(2.15), Inches(1.45), Inches(0.3),
                 scene, font_size=8, color=GREY, alignment=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.4), Inches(6.3), Inches(12), Inches(0.5),
             "⚠ Chapter breaks every 8-12 minutes for clean clipping | Email 5 must go out within 1 hour of show end",
             font_size=10, color=RED, bold=True)


def make_ghl_slide():
    """Slide 10: GHL Pipeline"""
    slide = add_slide("GHL PIPELINE & AUTOMATION", "6-Stage Guest Pipeline + 3-Stage Audience Pipeline")

    # Guest Pipeline
    add_box(slide, Inches(0.4), Inches(1.3), Inches(6.0), Inches(5.0), DARK)
    add_text(slide, Inches(0.6), Inches(1.4), Inches(5.6), Inches(0.4),
             "GUEST PIPELINE (6 Stages)", font_size=14, color=GOLD, bold=True)

    stages = [
        ("Prospect", "Identified potential guest"),
        ("Outreach Sent", "Initial invite email sent"),
        ("Shortlisted", "Responded positively"),
        ("Confirmed", "Booking confirmed"),
        ("Recorded", "Show complete"),
        ("Past Guest", "Clips sent, archived"),
    ]
    for i, (stage, desc) in enumerate(stages):
        color = [GOLD, MIDBLUE, GREEN, GREEN, RED, GREY][i]
        add_box(slide, Inches(0.6), Inches(1.9) + i * Inches(0.7), Inches(2.5), Inches(0.55),
                 color, stage, NAVY)
        add_text(slide, Inches(3.3), Inches(1.95) + i * Inches(0.7), Inches(2.8), Inches(0.5),
                 desc, font_size=9, color=WHITE)

    # Audience Pipeline
    add_box(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(3.0), DARK)
    add_text(slide, Inches(7.0), Inches(1.4), Inches(5.6), Inches(0.4),
             "AUDIENCE PIPELINE (3 Stages)", font_size=14, color=GOLD, bold=True)

    a_stages = [
        ("Show Subscriber", "Email captured; receiving announcements"),
        ("Engaged Viewer", "Commented, DMed, or attended 2+ shows"),
        ("CSQL", "Booked call from show content"),
    ]
    for i, (stage, desc) in enumerate(a_stages):
        add_box(slide, Inches(7.0), Inches(1.9) + i * Inches(0.7), Inches(2.5), Inches(0.55),
                 MIDBLUE, stage, WHITE)
        add_text(slide, Inches(9.7), Inches(1.95) + i * Inches(0.7), Inches(2.8), Inches(0.5),
                 desc, font_size=9, color=WHITE)


def make_postshow_slide():
    """Slide 11: Post-Show Sequence"""
    slide = add_slide("POST-SHOW SEQUENCE", "Emails 5, 6, 7 — Clips, replay, nurture")

    emails = [
        ("EMAIL 5 — Within 1 hour", "Guest thank-you + full recording",
         ["Thank you + recording link", "LinkedIn replay link",
          "Podcast upload notification"], RED),
        ("EMAIL 6 — 24 hours", "Audience replay + newsletter",
         ["Replay link for subscribers", "Newsletter sign-up CTA",
          "Next episode teaser"], MIDBLUE),
        ("EMAIL 7 — 24-36 hours", "Clip delivery",
         ["3 polished clips (MP4 9:16)", "Suggested captions",
          "Tag handles for host + show",
          "Download links (Google Drive)"], GREEN),
    ]

    for i, (title, purpose, bullets, color) in enumerate(emails):
        x = Inches(0.4) + i * Inches(4.2)
        add_box(slide, x, Inches(1.3), Inches(4.0), Inches(5.5), DARK)
        add_text(slide, x + Inches(0.1), Inches(1.4), Inches(3.8), Inches(0.3),
                 title, font_size=12, color=color, bold=True)
        add_text(slide, x + Inches(0.1), Inches(1.7), Inches(3.8), Inches(0.3),
                 purpose, font_size=9, color=GREY)
        for j, b in enumerate(bullets):
            add_text(slide, x + Inches(0.2), Inches(2.1) + j * Inches(0.5),
                     Inches(3.6), Inches(0.4), "• " + b, font_size=9, color=WHITE)


def make_skillmap_slide():
    """Slide 12: Skill Map"""
    slide = add_slide("SKILL MAP", "How the 5 Reveting skills connect")

    skills = [
        ("1. SHOW SETUP", "Master variable intake\nTier 1 + Tier 2 variables\nEmail trigger checklist", GOLD),
        ("2. EMAIL FLOWS", "6-sequence system\n11 email variants\nRe-engagement logic", MIDBLUE),
        ("3. CALENDAR FLOWS", "Booking pages\n6-touch reminders\nNo-show recovery", GREEN),
        ("4. GHL WORKFLOWS", "6-stage pipeline\n4 workflow maps\nContact schema", RED),
        ("5. STREAMYARD", "6-scene library\nGuest invite SOP\n45-min run-of-show", GOLD),
    ]

    for i, (title, desc, color) in enumerate(skills):
        y = Inches(1.3) + i * Inches(1.1)
        add_box(slide, Inches(0.4), y, Inches(12.5), Inches(0.95), DARK)
        add_box(slide, Inches(0.4), y, Inches(0.15), Inches(0.95), color)
        add_text(slide, Inches(0.7), y + Inches(0.05), Inches(3), Inches(0.4),
                 title, font_size=12, color=color, bold=True)
        add_text(slide, Inches(0.7), y + Inches(0.4), Inches(12), Inches(0.5),
                 desc, font_size=10, color=WHITE)


def make_summary_slide():
    """Slide 13: Summary / Next Steps"""
    slide = add_slide("SUMMARY", "Your show production checklist")

    add_box(slide, Inches(0.4), Inches(1.3), Inches(6.0), Inches(5.5), DARK)
    add_text(slide, Inches(0.6), Inches(1.4), Inches(5.6), Inches(0.4),
             "PRODUCTION CHECKLIST", font_size=14, color=GOLD, bold=True)

    checklist = [
        "☐ Tier 1 variables completed (show constants)",
        "☐ Guest booking confirmed in GHL",
        "☐ 6-touch reminder sequence activated",
        "☐ StreamYard scenes built and tested",
        "☐ LinkedIn Live event created",
        "☐ Multi-destination streaming configured",
        "☐ Email sequences loaded in GHL workflows",
        "☐ Post-production handoff SLA confirmed (24h clips)",
        "☐ Guest clips delivered within 36 hours",
        "☐ Episode uploaded to podcast channels",
    ]
    for i, item in enumerate(checklist):
        add_text(slide, Inches(0.7), Inches(1.9) + i * Inches(0.45),
                 Inches(5.4), Inches(0.4), item, font_size=10, color=WHITE)

    add_box(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), DARK)
    add_text(slide, Inches(7.0), Inches(1.4), Inches(5.6), Inches(0.4),
             "RESOURCES", font_size=14, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(7.1), Inches(1.9), Inches(5.4), Inches(4.0),
                    ["github.com/sitedesigner/golivestream",
                     "5 operational skills + templates",
                     "Sample emails and calendar events",
                     "PPTX production deck generator",
                     "Canva banner generator (coming soon)",
                     "GHL → Google Calendar sync (coming soon)"],
                    font_size=11)


def main():
    parser = argparse.ArgumentParser(description="Generate Reveting Show Operations PPTX")
    parser.add_argument("--output", default="Reveting-Show-Flow.pptx", help="Output filename")
    parser.add_argument("--show", default="The David Daily Show", help="Show name")
    parser.add_argument("--guest", default="", help="Guest name")
    args = parser.parse_args()

    global prs
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print(f"📊 Generating show operations deck...")
    print(f"   Show: {args.show}")
    if args.guest:
        print(f"   Guest: {args.guest}")

    make_title_slide()
    make_overview_slide()
    make_tier1_slide()
    make_guest_booking_slide()
    make_t14_slide()
    make_email1_slide()
    make_preshow_slide()
    make_streamyard_slide()
    make_live_slide()
    make_ghl_slide()
    make_postshow_slide()
    make_skillmap_slide()
    make_summary_slide()

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), args.output)
    prs.save(output_path)
    print(f"✅ Saved: {output_path}")
    print(f"   Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
